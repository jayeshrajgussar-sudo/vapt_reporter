from pptx import Presentation
from pptx.util import Inches, Pt
import os

# === USER INPUT ===
pptx_file = input("Enter path to the .pptx file (e.g., report.pptx): ").strip()
client_name = input("Enter Client Name: ").strip()
service_provider = input("Enter Service Provider Name: ").strip()
logo_path = input("Enter logo image file name (e.g., logo.png): ").strip()

# === CONFIGURATION ===
output_file = f"{os.path.splitext(pptx_file)[0]}_with_logo.pptx"
logo_width = Inches(1.5)
logo_top = Inches(0.1)
placeholder_mapping = {
    "client": client_name,
    "service provider": service_provider
}
# =====================

# Validate .pptx file
if not os.path.isfile(pptx_file):
    print(f"❌ Error: PPTX file '{pptx_file}' not found.")
    exit()

# Validate logo file
if not os.path.isfile(logo_path):
    print(f"❌ Error: Logo file '{logo_path}' not found.")
    exit()

# Load presentation
prs = Presentation(pptx_file)
slide_width = prs.slide_width

# Helper: Replace placeholders
def replace_text_placeholders(shape, mapping):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for key, val in mapping.items():
                    run.text = run.text.replace(f"[{key}]", val)
    elif shape.shape_type == 19:  # Table
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        for key, val in mapping.items():
                            run.text = run.text.replace(f"[{key}]", val)

# Process each slide
for slide_num, slide in enumerate(prs.slides, 1):
    # Replace placeholders
    for shape in slide.shapes:
        replace_text_placeholders(shape, placeholder_mapping)

    # Calculate logo position (top-right)
    logo_left = slide_width - logo_width - Inches(0.2)

    try:
        slide.shapes.add_picture(
            logo_path,
            logo_left,
            logo_top,
            width=logo_width
        )
    except Exception as e:
        print(f"❌ Failed to add logo to slide {slide_num}: {e}")

# Save result
prs.save(output_file)
print(f"✅ Presentation updated and saved as: {output_file}")
