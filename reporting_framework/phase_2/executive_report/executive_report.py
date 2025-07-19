import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

# === Import config ===
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../../")))
from config import CONFIG

# === Load config paths ===
conf = CONFIG["comparison_report"]
patched_file = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../", conf["patched_file"]))
unpatched_file = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../", conf["unpatched_file"]))
output_filename = conf["output"]

# === Output Path ===
base_output = os.environ.get("CUSTOM_OUTPUT_DIR", os.getcwd())
output_subdir = os.path.join(base_output, "phase_2", "executive_report")
os.makedirs(output_subdir, exist_ok=True)
output_pptx = os.path.join(output_subdir, output_filename)

# === Load and clean ===
def load_clean_df(file):
    df = pd.read_excel(file, engine="openpyxl")
    df.columns = df.columns.str.lower()
    df['nvt name'] = df['nvt name'].astype(str).str.strip()
    df['nvt_key'] = df['nvt name'].str.lower().str.replace('\xa0', ' ', regex=False)

    severity_col = next((col for col in df.columns if col.strip().lower() in ['severity', 'serverity', 'sev']), None)
    if severity_col:
        df['severity'] = df[severity_col].fillna('0').astype(str).str.strip()
    else:
        df['severity'] = '0'

    severity_map = {
        '0': 'Informational',
        '1': 'Low',
        '2': 'Medium',
        '3': 'High'
    }
    df['severity'] = df['severity'].apply(lambda x: severity_map.get(x, str(x).capitalize()))
    return df

patched_df = load_clean_df(patched_file)
unpatched_df = load_clean_df(unpatched_file)

# === Identify status sets ===
patched_set = set(patched_df['nvt_key'])
unpatched_set = set(unpatched_df['nvt_key'])

only_unpatched = unpatched_set - patched_set
both = patched_set & unpatched_set
only_patched = patched_set - unpatched_set

def deduplicate(df, keys, status):
    subset = df[df['nvt_key'].isin(keys)].copy()
    subset = subset.sort_values('severity').drop_duplicates('nvt_key', keep='first')
    subset['status'] = status
    return subset

open_new = deduplicate(unpatched_df, only_unpatched, 'Open')
open_both = deduplicate(patched_df, both, 'Open')
closed = deduplicate(patched_df, only_patched, 'Closed')

# === Combine and format ===
combined_df = pd.concat([open_new, open_both, closed], ignore_index=True)
combined_df = combined_df[['nvt name', 'severity', 'status']].copy()

severity_order = {'High': 1, 'Medium': 2, 'Low': 3, 'Informational': 4}
status_order = {'Open': 0, 'Closed': 1}
combined_df['severity_rank'] = combined_df['severity'].map(severity_order)
combined_df['status_rank'] = combined_df['status'].map(status_order)
combined_df = combined_df.sort_values(by=['severity_rank', 'status_rank']).reset_index(drop=True)
combined_df.drop(columns=['severity_rank', 'status_rank'], inplace=True)
combined_df.insert(0, 'S.No', range(1, len(combined_df) + 1))

# === Create Presentation ===
prs = Presentation()
prs.slide_width = Cm(21.59)
prs.slide_height = Cm(27.94)

def get_color(severity):
    return {
        "High": RGBColor(255, 0, 0),
        "Medium": RGBColor(255, 165, 0),
        "Low": RGBColor(0, 128, 0),
        "Informational": RGBColor(0, 0, 255)
    }.get(severity, RGBColor(128, 128, 128))

entries_per_slide = 18
cols = ['S.No', 'NVT Name', 'Severity', 'Status']

# === Generate Slides ===
for i in range(0, len(combined_df), entries_per_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(18), Cm(1.5))
    title_frame = title_box.text_frame
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "NVT Comparison Table"
    title_run.font.size = Pt(20)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)

    chunk = combined_df.iloc[i:i + entries_per_slide]
    rows = len(chunk) + 1
    table = slide.shapes.add_table(rows, 4, Cm(1.5), Cm(3), Cm(18), Cm(0.9) * rows).table

    col_widths = [Cm(2.0), Cm(10.5), Cm(3.0), Cm(3.5)]
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = width

    # Header Row
    for col_idx, col_name in enumerate(cols):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        for p in cell.text_frame.paragraphs:
            run = p.add_run()
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0, 0, 0)

    # Fill Data Rows
    for row_idx, (_, row) in enumerate(chunk.iterrows(), start=1):
        values = [row['S.No'], row['nvt name'], row['severity'], row['status']]
        severity_color = get_color(row['severity'])

        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            p = cell.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0, 0, 0)  # Static black text

            if col_idx == 3 and row['status'] == "Open":
                run.font.bold = True

            # Apply fill color only for Severity column
            if col_idx == 2:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                existing_fill = tcPr.find(qn('a:solidFill'))
                if existing_fill is not None:
                    tcPr.remove(existing_fill)
                solidFill = OxmlElement('a:solidFill')
                srgbClr = OxmlElement('a:srgbClr')
                srgbClr.set('val', {
                    "High": 'FF0000',
                    "Medium": 'FFA500',
                    "Low": '008000',
                    "Informational": '0000FF'
                }.get(row['severity'], '808080'))
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

# === Save File ===
prs.save(output_pptx)
print(f"Report saved successfully at: {output_pptx}")
