import os
import copy
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor

# === Allow standalone execution ===
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../../")))
from config import CONFIG

# === Load config values ===
conf = CONFIG["detailed_nvt_report"]
excel_file = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../", conf["excel_file"]))
template_file = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../", conf["template"]))
output_filename = conf["output"]

# === Set output path ===
base_output = os.environ.get("CUSTOM_OUTPUT_DIR", os.getcwd())
output_dir = os.path.join(base_output, "phase_1", "detailed_nvt_reports")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, output_filename)

# === Load Excel ===
df = pd.read_excel(excel_file)
if df.empty:
    print("Error: Excel file is empty. Please add data.")
    exit()

# === Clean & transform ===
columns_needed = [
    'IP', 'Port', 'Port Protocol', 'Severity', 'Solution Type',
    'NVT Name', 'Summary', 'Solution', 'Vulnerability Insight', 'CVEs'
]

# Fill missing columns
for col in columns_needed:
    if col not in df.columns:
        df[col] = ''

# Fill NaN and strip whitespace
for col in columns_needed:
    df[col] = df[col].fillna('').astype(str).str.strip()

def normalize_severity(x):
    severity_map = {
        0: 'Informational',
        1: 'Low',
        2: 'Medium',
        3: 'High'
    }
    try:
        if isinstance(x, str):
            x_clean = x.strip().lower()
            if x_clean in ['informational', 'low', 'medium', 'high']:
                return x_clean.capitalize()
            x_numeric = float(x_clean)
            return severity_map.get(int(x_numeric), 'Informational')
        elif isinstance(x, (int, float)):
            return severity_map.get(int(x), 'Informational')
    except:
        pass
    return 'Informational'


df['Severity'] = df['Severity'].fillna(0)
df['Severity'] = df['Severity'].apply(normalize_severity)


# Normalize port
df['Port'] = df['Port'].apply(lambda x: str(int(float(x))) if str(x).replace('.', '', 1).isdigit() else str(x).strip())

# Grouping key
df['ip_combo'] = df['IP'] + '/' + df['Port'] + '/' + df['Port Protocol']

# Rank severities
severity_order = {
    'High': 1,
    'Medium': 2,
    'Low': 3,
    'Informational': 4
}
df['severity_rank'] = df['Severity'].map(lambda x: severity_order.get(x, 5))

# === Grouping ===
grouped_df = df.groupby(
    ['Severity', 'Solution Type', 'NVT Name', 'Summary', 'Solution', 'Vulnerability Insight', 'CVEs']
).agg({
    'IP': lambda x: ', '.join(sorted(set(x))),
    'Port': lambda x: ', '.join(sorted(set(x))),
    'Port Protocol': lambda x: ', '.join(sorted(set(x))),
    'severity_rank': 'first'
}).reset_index()

grouped_df = grouped_df.sort_values(by='severity_rank').drop(columns='severity_rank')
grouped_df.insert(0, 'no', range(1, len(grouped_df) + 1))

# === Helpers ===
def clean_cell(value, is_multiline=False):
    if pd.isna(value):
        return ''
    text = str(value).strip().replace('\r\n', '\n').replace('\r', '\n').replace('\n\n', '\n')
    if is_multiline:
        text = text.replace('\n', '. ').replace('..', '.')
        if not text.endswith('.'):
            text += '.'
    return text

def get_severity_color(severity):
    return {
        "High": RGBColor(255, 0, 0),
        "Medium": RGBColor(255, 165, 0),
        "Low": RGBColor(0, 128, 0),
        "Informational": RGBColor(0, 0, 255)
    }.get(severity, RGBColor(128, 128, 128))  # default gray

def replace_text_in_shape(shape, mapping):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full_text = ''.join(run.text for run in para.runs)
            para.clear()
            run = para.add_run()

            if "[Severity]" in full_text:
                sev = mapping.get('Severity', 'Informational')
                color = get_severity_color(sev)
                run.text = full_text.replace("[Severity]", sev)
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(255, 255, 255)

                shape.fill.solid()
                shape.fill.fore_color.rgb = color
            else:
                for key, val in mapping.items():
                    full_text = full_text.replace(f"[{key}]", val)
                run.text = full_text
                run.font.size = Pt(11)

    elif shape.shape_type == 19 and hasattr(shape, "table"):  # Table
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    full_text = ''.join(run.text for run in para.runs)
                    para.clear()
                    run = para.add_run()

                    if "[Severity]" in full_text:
                        sev = mapping.get("Severity", "Informational")
                        color = get_severity_color(sev)
                        run.text = full_text.replace("[Severity]", sev)
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(255, 255, 255)

                        cell.fill.solid()
                        cell.fill.fore_color.rgb = color
                    else:
                        for key, val in mapping.items():
                            full_text = full_text.replace(f"[{key}]", val)
                        run.text = full_text
                        run.font.size = Pt(11)

# === Build output presentation ===
template_prs = Presentation(template_file)
output_prs = Presentation()
output_prs.slide_width = Cm(21.59)
output_prs.slide_height = Cm(27.94)
template_slide = template_prs.slides[0]

for _, row in grouped_df.iterrows():
    mapping = {
        "no": str(row["no"]),
        "IP": clean_cell(row["IP"]),
        "Port": clean_cell(row["Port"]),
        "Port Protocol": clean_cell(row["Port Protocol"]),
        "Severity": clean_cell(row["Severity"]),
        "Solution Type": clean_cell(row["Solution Type"]),
        "NVT Name": clean_cell(row["NVT Name"]),
        "Summary": clean_cell(row["Summary"], is_multiline=True),
        "CVEs": clean_cell(row["CVEs"]),
        "Solution": clean_cell(row["Solution"], is_multiline=True),
        "Vulnerability Insight": clean_cell(row["Vulnerability Insight"], is_multiline=True)
    }

    slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
    for shape in template_slide.shapes:
        if not shape.is_placeholder:
            slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), 'p:extLst')
    for shape in slide.shapes:
        replace_text_in_shape(shape, mapping)

output_prs.save(output_path)
print(f"Detailed NVT report saved at: {output_path}")
