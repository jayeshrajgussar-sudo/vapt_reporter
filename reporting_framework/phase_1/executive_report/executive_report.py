import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

# === Import config ===
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../../")))
from config import CONFIG

# === Load Configuration ===
conf = CONFIG["executive_report"]
excel_file = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../", conf["excel_file"]))
output_filename = conf["output"]

# === Output Path ===
base_output = os.environ.get("CUSTOM_OUTPUT_DIR", os.getcwd())
output_dir = os.path.join(base_output, "phase_1", "executive_reports")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, output_filename)

try:
    print("Checking Excel file...")
    if not os.path.exists(excel_file):
        raise FileNotFoundError(f"Excel file not found: {excel_file}")
    if os.path.getsize(excel_file) == 0:
        raise ValueError("Excel file is empty")

    print(f"Reading: {excel_file}")
    df = pd.read_excel(excel_file, engine="openpyxl")

    if df.empty:
        raise ValueError("Excel file contains no data")

    print("Cleaning and preparing data...")
    df['NVT Name'] = df['NVT Name'].astype(str).str.strip()

    # Replace NaN or missing values in Severity
    df['Severity'] = df['Severity'].fillna(0)

    # Normalize severity
    def normalize_severity(x):
        severity_map = {
            0: 'Informational',
            1: 'Low',
            2: 'Medium',
            3: 'High',
            '0': 'Informational',
            '1': 'Low',
            '2': 'Medium',
            '3': 'High',
            'informational': 'Informational',
            'low': 'Low',
            'medium': 'Medium',
            'high': 'High'
        }
        return severity_map.get(x, severity_map.get(str(x).strip().lower(), 'Informational'))

    df['Severity'] = df['Severity'].apply(normalize_severity)

    df['status'] = 'Open'
    grouped = df[['NVT Name', 'Severity', 'status']].drop_duplicates()

    severity_order = {
        'High': 1,
        'Medium': 2,
        'Low': 3,
        'Informational': 4
    }
    grouped['severity_rank'] = grouped['Severity'].map(lambda x: severity_order.get(x, 5))
    grouped = grouped.sort_values(by='severity_rank').reset_index(drop=True)
    grouped.insert(0, 'S.No', range(1, len(grouped) + 1))

    print("Generating presentation...")
    prs = Presentation()
    prs.slide_width = Cm(21.59)
    prs.slide_height = Cm(27.94)

    layout = prs.slide_layouts[6]
    max_rows_per_slide = 19
    col_widths = [Cm(2), Cm(10), Cm(4), Cm(3.5)]
    headers = ['S.No', 'Name Of Vulnerability', 'Severity', 'Status']

    def get_severity_color(severity):
        return {
            "High": RGBColor(255, 0, 0),          # Red
            "Medium": RGBColor(255, 165, 0),      # Orange
            "Low": RGBColor(102, 255, 102),       # Light Green
            "Informational": RGBColor(0, 0, 255)  # Blue
        }.get(severity, RGBColor(0, 0, 255))       # Default to blue

    chunks = [grouped.iloc[i:i + max_rows_per_slide] for i in range(0, len(grouped), max_rows_per_slide)]

    for chunk in chunks:
        slide = prs.slides.add_slide(layout)
        rows, cols = len(chunk) + 1, 4
        table = slide.shapes.add_table(rows, cols, Cm(1), Cm(2), Cm(19.5), Cm(1.2) * rows).table

        for i, width in enumerate(col_widths):
            table.columns[i].width = width

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)

        # Data rows
        for row_idx, (_, row) in enumerate(chunk.iterrows(), start=1):
            for col_idx, key in enumerate(['S.No', 'NVT Name', 'Severity', 'status']):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(row[key])
                if key == 'Severity':
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = get_severity_color(row[key])
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(11)

    print("Saving file...")
    prs.save(output_path)
    print(f"Report generated successfully at: {output_path}")

except Exception as e:
    print(f"Error: {e}")
