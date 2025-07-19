import os
import json
import re
import sys
import shutil
from pptx import Presentation

sys.stdout.reconfigure(line_buffering=True)

# === Load config ===
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../../")))
from config import CONFIG

# === Placeholder setup ===
base_placeholders = [
    "service_name", "client", "service_provider", "report_release_date",
    "type_of_audit", "type_of_audit_report", "period", "document_title",
    "document_id", "prepared_by", "reviewed_by", "approved_by", "released_by",
    "release_date", "number_of_locations", "parent_client_name", "asset/range_location",
    "pre_audit_prep", "activities_dates", "first_level_report", "second_level_report",
    "client_location"
]

alias_map = {
    "clinet": "client",
    "server_name": "asset/range_location",
    "server_location": "asset/range_location",
    " client_location ": "client_location"
}

def all_placeholder_forms(k):
    return [f"[{k}]", f"[ {k} ]", f"{{{k}}}", f"({k})", f"([{k}])"]

def normalize_text(text):
    for alias, actual in alias_map.items():
        pattern = rf"(\[|\(|\{{)\s*{re.escape(alias)}\s*(\]|\)|\}})"
        text = re.sub(pattern, f"[{actual}]", text, flags=re.IGNORECASE)
    return text

def replace_placeholders(text, replacements):
    text = normalize_text(text)
    for key, val in replacements.items():
        for form in all_placeholder_forms(key):
            text = re.sub(re.escape(form), val, text, flags=re.IGNORECASE)
    return text

def replace_text_frame(text_frame, replacements):
    full_text = "".join(run.text for para in text_frame.paragraphs for run in para.runs)
    replaced = replace_placeholders(full_text, replacements)
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        text_frame.paragraphs[0].runs[0].text = replaced
        for para in text_frame.paragraphs:
            for run in para.runs[1:]:
                run.text = ""

# === Load inputs.json ===
input_path = os.path.join(os.path.dirname(__file__), "inputs.json")
if not os.path.exists(input_path):
    print("Error: inputs.json not found.")
    sys.exit(1)

with open(input_path, "r") as f:
    replacements = json.load(f)

# === Load template ===
template_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../", CONFIG["front_page"]["input_pptx"]))
if not os.path.exists(template_path):
    print(f"Error: Template not found: {template_path}")
    sys.exit(1)

prs = Presentation(template_path)

# === Replace placeholders in slides ===
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text_frame(shape.text_frame, replacements)
        elif shape.shape_type == 19:  # Table
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_text_frame(cell.text_frame, replacements)

# === Save to: /outputs/<custom_folder>/phase_1/front/output_file.pptx ===
base_output = os.environ.get("CUSTOM_OUTPUT_DIR", os.getcwd())
output_dir_1 = os.path.join(base_output, "phase_1", "front")
os.makedirs(output_dir_1, exist_ok=True)

output_filename = CONFIG["front_page"]["output_pptx"]
output_path_1 = os.path.join(output_dir_1, output_filename)
prs.save(output_path_1)

# Also copy to phase_2/front if needed
output_dir_2 = os.path.join(base_output, "phase_2", "front")
os.makedirs(output_dir_2, exist_ok=True)
shutil.copyfile(output_path_1, os.path.join(output_dir_2, output_filename))

print(f"Saved front slide:\n{output_path_1}")

# === Cleanup: Delete inputs.json ===
try:
    os.remove(input_path)
    print(" Cleaned up: inputs.json deleted.")
except Exception as e:
    print(f" Failed to delete inputs.json: {e}")
